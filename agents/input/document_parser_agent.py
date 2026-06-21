import os
import io
import base64
from typing import Any, Dict, List, Optional
from pathlib import Path
import asyncio

from core.base_agent import BaseAgent, TaskResult
from core.llm_client import LLMClient, LLMProvider

class DocumentParserAgent(BaseAgent):
    def __init__(self):
        super().__init__(
            name="DocumentParserAgent",
            description="解析多模态档案文件，支持文本、图片、视频、音频、3D扫描件"
        )
        self.supported_formats = {
            "text": [".txt", ".md", ".csv", ".json", ".xml"],
            "document": [".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"],
            "image": [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff", ".webp"],
            "video": [".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv"],
            "audio": [".mp3", ".wav", ".flac", ".aac", ".ogg", ".m4a"],
            "model3d": [".obj", ".fbx", ".gltf", ".glb", ".ply", ".stl"]
        }
        self.llm_client = LLMClient(LLMProvider.QWEN)
    
    async def execute(self, task_input: Dict[str, Any]) -> TaskResult:
        files = task_input.get("files", [])
        if not files:
            return TaskResult(success=False, error="No files provided")
        
        parsed_results = []
        for file_info in files:
            try:
                result = await self._parse_file(file_info)
                parsed_results.append(result)
            except Exception as e:
                parsed_results.append({
                    "file": file_info.get("path", "unknown"),
                    "error": str(e),
                    "type": "error"
                })
        
        return TaskResult(
            success=True,
            data={
                "parsed_files": parsed_results,
                "total_count": len(files),
                "success_count": sum(1 for r in parsed_results if "error" not in r)
            },
            metadata={"agent": self.name}
        )
    
    async def _parse_file(self, file_info: Dict[str, Any]) -> Dict[str, Any]:
        file_path = file_info.get("path")
        file_type = file_info.get("type", "auto")
        
        if not file_path or not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        path = Path(file_path)
        suffix = path.suffix.lower()
        
        # 自动检测类型
        if file_type == "auto":
            file_type = self._detect_file_type(suffix)
        
        parser_methods = {
            "text": self._parse_text,
            "document": self._parse_document,
            "image": self._parse_image,
            "video": self._parse_video,
            "audio": self._parse_audio,
            "model3d": self._parse_3d_model
        }
        
        parser = parser_methods.get(file_type)
        if not parser:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return await parser(file_path, file_info)
    
    def _detect_file_type(self, suffix: str) -> str:
        for file_type, extensions in self.supported_formats.items():
            if suffix in extensions:
                return file_type
        return "unknown"
    
    async def _parse_text(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        return {
            "file": file_path,
            "type": "text",
            "content": content,
            "length": len(content),
            "encoding": "utf-8"
        }
    
    async def _parse_document(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        suffix = Path(file_path).suffix.lower()
        
        if suffix == ".pdf":
            return await self._parse_pdf(file_path)
        elif suffix in [".doc", ".docx"]:
            return await self._parse_word(file_path)
        elif suffix in [".ppt", ".pptx"]:
            return await self._parse_ppt(file_path)
        elif suffix in [".xls", ".xlsx"]:
            return await self._parse_excel(file_path)
        else:
            raise ValueError(f"Unsupported document format: {suffix}")
    
    async def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        try:
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                
                return {
                    "file": file_path,
                    "type": "pdf",
                    "content": text,
                    "pages": len(reader.pages),
                    "metadata": {
                        "title": reader.metadata.get("/Title", ""),
                        "author": reader.metadata.get("/Author", ""),
                        "creation_date": str(reader.metadata.get("/CreationDate", ""))
                    }
                }
        except ImportError:
            # 如果没有PyPDF2，返回文件信息
            return {
                "file": file_path,
                "type": "pdf",
                "content": "PDF parsing requires PyPDF2 library",
                "note": "Install with: pip install PyPDF2"
            }
    
    async def _parse_word(self, file_path: str) -> Dict[str, Any]:
        try:
            from docx import Document
            import os as _os
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            # 提取 DOCX 内嵌图片
            embedded_images = []
            img_dir = _os.path.join(_os.path.dirname(file_path), "_docx_images")
            _os.makedirs(img_dir, exist_ok=True)
            doc_name = _os.path.splitext(_os.path.basename(file_path))[0]

            for i, rel in enumerate(doc.part.rels.values()):
                if "image" in rel.reltype:
                    img = rel.target_part
                    ext = _os.path.splitext(img.partname)[1] or ".png"
                    img_filename = f"{doc_name}_img{i}{ext}"
                    img_path = _os.path.join(img_dir, img_filename)
                    with open(img_path, "wb") as f:
                        f.write(img.blob)
                    embedded_images.append({
                        "local_path": img_path,
                        "filename": img_filename,
                        "size_bytes": len(img.blob),
                    })

            # 提取表格内容
            table_data = []
            for t_idx, table in enumerate(doc.tables):
                rows = []
                for row in table.rows:
                    cells = [cell.text for cell in row.cells]
                    rows.append(cells)
                table_data.append({"index": t_idx, "rows": rows, "row_count": len(rows)})

            return {
                "file": file_path,
                "type": "word",
                "content": text,
                "paragraphs": len(doc.paragraphs),
                "tables": len(doc.tables),
                "table_data": table_data,
                "embedded_images": embedded_images,
                "embedded_image_count": len(embedded_images),
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "word",
                "content": "Word parsing requires python-docx library",
            }
    
    async def _parse_ppt(self, file_path: str) -> Dict[str, Any]:
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_content = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                slides_content.append({
                    "slide_number": i + 1,
                    "content": "\n".join(slide_text)
                })
            
            return {
                "file": file_path,
                "type": "ppt",
                "slides": slides_content,
                "slide_count": len(prs.slides)
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "ppt",
                "content": "PPT parsing requires python-pptx library",
                "note": "Install with: pip install python-pptx"
            }
    
    async def _parse_excel(self, file_path: str) -> Dict[str, Any]:
        try:
            import pandas as pd
            df = pd.read_excel(file_path)
            
            return {
                "file": file_path,
                "type": "excel",
                "content": df.to_string(),
                "columns": list(df.columns),
                "rows": len(df),
                "data": df.to_dict(orient='records')
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "excel",
                "content": "Excel parsing requires pandas and openpyxl",
                "note": "Install with: pip install pandas openpyxl"
            }
    
    async def _parse_image(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        try:
            from PIL import Image
            import base64
            
            img = Image.open(file_path)
            
            # 转换为base64用于AI分析
            with open(file_path, "rb") as f:
                img_base64 = base64.b64encode(f.read()).decode('utf-8')
            
            # 使用多模态LLM分析图像内容
            image_analysis = await self._analyze_image_with_llm(img_base64, file_path)
            
            return {
                "file": file_path,
                "type": "image",
                "format": img.format,
                "size": img.size,
                "mode": img.mode,
                "base64": img_base64[:100] + "...",  # 截断显示
                "analysis": image_analysis
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "image",
                "note": "Image analysis requires PIL/Pillow",
                "install": "pip install Pillow"
            }
    
    async def _analyze_image_with_llm(self, img_base64: str, file_path: str) -> Dict[str, Any]:
        try:
            messages = [
                {
                    "role": "system",
                    "content": "你是一个档案图像分析专家。请详细描述这张档案图片的内容，包括文字、人物、场景、年代特征等。"
                },
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": "请分析这张档案图片，提取所有可见信息："
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{img_base64}"
                            }
                        }
                    ]
                }
            ]
            
            response = await self.llm_client.chat_completion(messages, max_tokens=2000)
            analysis_text = response["choices"][0]["message"]["content"]
            
            return {
                "description": analysis_text,
                "method": "multimodal_llm"
            }
        except Exception as e:
            return {
                "description": f"图像分析失败: {str(e)}",
                "method": "error"
            }
    
    async def _parse_video(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        try:
            import cv2
            cap = cv2.VideoCapture(file_path)
            
            fps = cap.get(cv2.CAP_PROP_FPS)
            frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
            height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
            duration = frame_count / fps if fps > 0 else 0
            
            # 提取关键帧
            keyframes = []
            frame_interval = max(1, frame_count // 10)  # 提取10帧
            
            for i in range(0, frame_count, frame_interval):
                cap.set(cv2.CAP_PROP_POS_FRAMES, i)
                ret, frame = cap.read()
                if ret:
                    # 保存关键帧
                    keyframe_path = f"{file_path}_frame_{i}.jpg"
                    cv2.imwrite(keyframe_path, frame)
                    keyframes.append({
                        "frame_number": i,
                        "timestamp": i / fps if fps > 0 else 0,
                        "path": keyframe_path
                    })
            
            cap.release()
            
            return {
                "file": file_path,
                "type": "video",
                "format": Path(file_path).suffix,
                "fps": fps,
                "frame_count": frame_count,
                "resolution": f"{width}x{height}",
                "duration": duration,
                "keyframes": keyframes
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "video",
                "note": "Video parsing requires opencv-python",
                "install": "pip install opencv-python"
            }
    
    async def _parse_audio(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        """音频解析 - 使用国产语音识别服务"""
        try:
            # 优先使用阿里云语音识别 (国产)
            return await self._parse_audio_with_ali(file_path)
        except Exception:
            # 降级：返回音频基本信息
            return await self._parse_audio_basic(file_path)

    async def _parse_audio_with_ali(self, file_path: str) -> Dict[str, Any]:
        """使用阿里云语音识别 (国产替代方案)"""
        try:
            import subprocess
            # 检查是否安装了 funasr (国产语音识别框架)
            result = subprocess.run(
                ["python", "-c", "import funasr; print('ok')"],
                capture_output=True, text=True
            )
            if result.returncode != 0:
                raise ImportError("funasr not installed")
            
            return {
                "file": file_path,
                "type": "audio",
                "note": "检测到 funasr (国产语音识别框架)",
                "transcription": "请使用 funasr 进行离线语音识别",
                "framework": "funasr (阿里巴巴达摩院)",
                "install": "pip install funasr modelscope"
            }
        except Exception as e:
            raise e

    async def _parse_audio_basic(self, file_path: str) -> Dict[str, Any]:
        """返回音频基本信息"""
        import os
        file_size = os.path.getsize(file_path)
        
        # 尝试获取音频时长
        duration = 0
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(file_path)
            duration = len(audio) / 1000  # 转换为秒
        except ImportError:
            pass
        
        return {
            "file": file_path,
            "type": "audio",
            "size_bytes": file_size,
            "duration_seconds": duration,
            "note": "音频文件已识别，语音识别需要配置国产ASR服务",
            "supported_asr": [
                "阿里云语音识别 (nls-python-sdk)",
                "百度语音识别 (baidu-aip)",
                "讯飞语音识别 (xfyun)",
                "funasr (阿里巴巴达摩院开源)"
            ]
        }
    
    async def _parse_3d_model(self, file_path: str, file_info: Dict) -> Dict[str, Any]:
        try:
            import trimesh
            mesh = trimesh.load(file_path)
            
            return {
                "file": file_path,
                "type": "3d_model",
                "format": Path(file_path).suffix,
                "vertices": len(mesh.vertices),
                "faces": len(mesh.faces),
                "bounds": mesh.bounds.tolist(),
                "center": mesh.centroid.tolist(),
                "note": "3D model loaded successfully"
            }
        except ImportError:
            return {
                "file": file_path,
                "type": "3d_model",
                "note": "3D parsing requires trimesh",
                "install": "pip install trimesh"
            }
